import glob
import cv2
import sys
import os
import ipdb
from pptx import Presentation
from pptx.util import Inches, Pt
os.system('cp -rf ../../code/python/eis_core ./')
os.system('cp ../../code/python/config_path.py ./')
from eis_core import eis_core_config_loader
import config_path



if len(sys.argv)!=5:
    print("Please enter: gen_report.py dataset_dir experiment_prefix ppt_title ppt_subtitle")
    exit()
else:
    dataset_dir = sys.argv[1]
    experiment_prefix = sys.argv[2]
    ppt_title = sys.argv[3]
    ppt_subtitle  = sys.argv[4]


# load debug/eis param
debug_param = eis_core_config_loader.load_debug_config(config_path.debug_config_path)
eis_param = eis_core_config_loader.load_eis_config(config_path.eis_config_path)



ppt_path = dataset_dir + '/' + ppt_title
experiment_path = dataset_dir+'/'+experiment_prefix+'*'
experiment_dirs = glob.glob(experiment_path) 


print('Generate Report:%s'%ppt_path)

#result log path
if eis_param.use_nonlinear_filter:
    result_dir = '/result' + '_nonlinear_alpha_min_%.2f_max_%.2f_beta_%.2f_gamma_%.2f_inner_ratio_%.1f_lookahead_%d_crop_ratio_%.1f'\
    %(eis_param.alpha_min,eis_param.alpha_max,eis_param.beta,eis_param.gamma,eis_param.inner_padding_ratio,eis_param.lookahead_no,eis_param.crop_ratio) 
else:
    result_dir = '/result' + '_const_lpf_alpha_%.2f_inner_ratio_%.1f_lookahead_%d_crop_ratio_%.1f'%(eis_param.alpha,eis_param.inner_padding_ratio,eis_param.lookahead_no,eis_param.crop_ratio) 

# create slide ------
prs = Presentation()
#------------slide layout--------------------------------
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]
title_only_slide_layout = prs.slide_layouts[5]

# -------------add title page---------------------------
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
# set slide title
title.text = ppt_title.split('.')[0]
subtitle.text = ppt_subtitle

# -------------add algo parameters page---------------------------
experiment_slide = prs.slides.add_slide(title_only_slide_layout)

if eis_param.use_nonlinear_filter:
    
    shapes = experiment_slide.shapes
    title_shape = shapes.title
    title_shape.text = 'non-linear filter'
    rows, cols, left, top, width, height = 2, 7, Inches(0.5), Inches(1.5), Inches(9), Inches(1)
    table = experiment_slide.shapes.add_table(rows, cols, left, top, width, height).table 
    table.cell(0, 0).text = 'alpha_min'
    table.cell(0, 1).text = 'alpha_max'
    table.cell(0, 2).text = 'beta'
    table.cell(0, 3).text = 'gamma'
    table.cell(0, 4).text = 'inner_ratio'
    table.cell(0, 5).text = 'lookahead'
    table.cell(0, 6).text = 'crop_ratio'

    table.cell(1, 0).text = '%.2f'%eis_param.alpha_min
    table.cell(1, 1).text = '%.2f'%eis_param.alpha_max
    table.cell(1, 2).text = '%.2f'%eis_param.beta
    table.cell(1, 3).text = '%.2f'%eis_param.gamma
    table.cell(1, 4).text = '%.2f'%eis_param.inner_padding_ratio
    table.cell(1, 5).text = '%.2f'%eis_param.lookahead_no
    table.cell(1, 6).text = '%.2f'%eis_param.crop_ratio

else:
    
    shapes = experiment_slide.shapes
    title_shape = shapes.title
    title_shape.text = 'constant filter'
    rows, cols, left, top, width, height = 2, 3, Inches(1), Inches(2.5), Inches(4.5), Inches(1)
    table = experiment_slide.shapes.add_table(rows, cols, left, top, width, height).table 
    table.cell(0, 0).text = 'alpha'
    table.cell(0, 1).text = 'inner_ratio'
    table.cell(0, 2).text = 'look_ahead'
    table.cell(1, 0).text = '%.2f'%eis_param.alpha
    table.cell(1, 1).text = '%.2f'%eis_param.inner_padding_ratio
    table.cell(1, 2).text = '%.2f'%eis_param.lookahead_no
    

#---------------- add experiment slide------------------------------


for i in range(len(experiment_dirs)):
    experiment_title = 'dataset_'+ experiment_dirs[i].split('/')[-1]
    experiment_dir = experiment_dirs[i] + '/' + result_dir
    print('generate dataset %s experiment report..'%experiment_dirs[i].split('/')[-1])
    # add video
    experiment_slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = experiment_slide.shapes
    title_shape = shapes.title
    title_shape.text = experiment_title
    video_path = glob.glob(experiment_dir+'/compare*.mp4')[0]
    video = cv2.VideoCapture(video_path)
    for j in range(15):
        ret, frame = video.read()
    cv2.imwrite('frame.png',frame)
    video.release()
    left, top, width, height = Inches(0), Inches(2), Inches(10), Inches(2.81) 
    experiment_slide.shapes.add_movie(video_path, left, top, width, height,poster_frame_image='frame.png')


    # add boundary control
    experiment_slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = experiment_slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = experiment_title
    tf = body_shape.text_frame
    tf.text = 'Boundary Control'
    img_path = experiment_dir + '/boundary_control.png'
    left, top, width, height = Inches(1), Inches(2.5), Inches(9), Inches(5)  
    pic = experiment_slide.shapes.add_picture(img_path, left, top, width, height)  

    # add path analysis
    experiment_slide = prs.slides.add_slide(title_slide_layout)
    img_path = experiment_dir + '/path_analysis.png'
    left, top, width, height = Inches(2), Inches(-0.5), Inches(6), Inches(8.5)  
    pic = experiment_slide.shapes.add_picture(img_path, left, top, width, height)


    # add path x
    experiment_slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = experiment_slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = experiment_title
    tf = body_shape.text_frame
    tf.text = 'camera axis_x'
    img_path = experiment_dir + '/cam_orien_x.png'
    left, top, width, height = Inches(1), Inches(2.5), Inches(9), Inches(5)  
    pic = experiment_slide.shapes.add_picture(img_path, left, top, width, height)  

    # add path y
    experiment_slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = experiment_slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = experiment_title
    tf = body_shape.text_frame
    tf.text = 'camera axis_y'
    img_path = experiment_dir + '/cam_orien_y.png'
    left, top, width, height = Inches(1), Inches(2.5), Inches(9), Inches(5)  
    pic = experiment_slide.shapes.add_picture(img_path, left, top, width, height)  

    # add path z
    experiment_slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = experiment_slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = experiment_title
    tf = body_shape.text_frame
    tf.text = 'camera axis_z'
    img_path = experiment_dir + '/cam_orien_z.png'
    left, top, width, height = Inches(1), Inches(2.5), Inches(9), Inches(5)  
    pic = experiment_slide.shapes.add_picture(img_path, left, top, width, height)  


print('save power point.....') 
prs.save(ppt_path)
os.system('rm -rf eis_core')
os.system('rm -rf config_path.py')
os.system('rm frame.png')